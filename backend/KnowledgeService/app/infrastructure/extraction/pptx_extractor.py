from __future__ import annotations

import io

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.application.ports.extractor import CorruptFileError, Extractor
from app.domain.extraction.extracted_image import ExtractedImage
from app.domain.extraction.extraction_result import ExtractionResult
from app.domain.extraction.text_block import TextBlock, TextBlockKind
from app.infrastructure.extraction._support import (
    PPTX_CONTENT_TYPES,
    PPTX_EXTENSIONS,
    OrderCounter,
    logger,
    normalize_content_type,
    read_image_dimensions,
    sha256_hex,
)


class PptxExtractor(Extractor):
    """Extractor for PowerPoint (.pptx) files, backed by python-pptx.

    Per slide (1-based): the title placeholder becomes a TITLE block, other text
    frames become PARAGRAPH blocks, tables become TABLE blocks, and speaker notes
    become a NOTES block. Picture shapes are inventoried with the slide number.
    """

    content_types = PPTX_CONTENT_TYPES
    extensions = PPTX_EXTENSIONS

    def supports(self, content_type: str) -> bool:
        return normalize_content_type(content_type) in self.content_types

    def extract(self, file_bytes: bytes, file_name: str) -> ExtractionResult:
        try:
            presentation = Presentation(io.BytesIO(file_bytes))
        except Exception as exc:  # PackageNotFoundError / bad zip
            raise CorruptFileError(
                f"Could not open PPTX {file_name!r}: {exc}"
            ) from exc

        order = OrderCounter()
        blocks: list[TextBlock] = []
        images: list[ExtractedImage] = []

        for slide_index, slide in enumerate(presentation.slides):
            slide_number = slide_index + 1
            title_shape = slide.shapes.title

            for shape in slide.shapes:
                if shape.has_table:
                    blocks.append(
                        self._table_block(shape, slide_number, order)
                    )
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = self._image(shape, slide_number, order)
                    if image is not None:
                        images.append(image)
                elif shape.has_text_frame:
                    block = self._text_block(
                        shape, title_shape, slide_number, order
                    )
                    if block is not None:
                        blocks.append(block)

            notes = self._notes_block(slide, slide_number, order)
            if notes is not None:
                blocks.append(notes)

        result = ExtractionResult(
            source_format="pptx",
            blocks=blocks,
            images=images,
            page_count=len(presentation.slides),
        )
        _log_summary(result)
        return result

    @staticmethod
    def _text_block(
        shape, title_shape, slide_number: int, order: OrderCounter
    ) -> TextBlock | None:
        text = shape.text_frame.text.strip()
        if not text:
            return None
        is_title = title_shape is not None and shape.shape_id == title_shape.shape_id
        return TextBlock(
            order=order.next(),
            text=text,
            kind=TextBlockKind.TITLE if is_title else TextBlockKind.PARAGRAPH,
            slide=slide_number,
        )

    @staticmethod
    def _table_block(shape, slide_number: int, order: OrderCounter) -> TextBlock:
        rows: list[str] = []
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                rows.append(" | ".join(cells))
        return TextBlock(
            order=order.next(),
            text="\n".join(rows),
            kind=TextBlockKind.TABLE,
            slide=slide_number,
        )

    @staticmethod
    def _notes_block(slide, slide_number: int, order: OrderCounter) -> TextBlock | None:
        if not slide.has_notes_slide:
            return None
        notes_frame = slide.notes_slide.notes_text_frame
        if notes_frame is None:
            return None
        text = notes_frame.text.strip()
        if not text:
            return None
        return TextBlock(
            order=order.next(),
            text=text,
            kind=TextBlockKind.NOTES,
            slide=slide_number,
        )

    @staticmethod
    def _image(shape, slide_number: int, order: OrderCounter) -> ExtractedImage | None:
        picture = shape.image
        data = picture.blob
        ext = (picture.ext or "").lower()
        dims = read_image_dimensions(data)
        if dims is None:
            logger.warning(
                "Skipping unreadable picture on PPTX slide %d.", slide_number
            )
            return None
        width, height, decoded_ext = dims
        return ExtractedImage(
            order=order.next(),
            data=data,
            ext=decoded_ext or ext,
            width=width,
            height=height,
            sha256=sha256_hex(data),
            slide=slide_number,
        )


def _log_summary(result: ExtractionResult) -> None:
    logger.info(
        "Extracted pptx: %d blocks, %d images, %d slide(s)",
        len(result.blocks),
        len(result.images),
        result.page_count or 0,
    )
