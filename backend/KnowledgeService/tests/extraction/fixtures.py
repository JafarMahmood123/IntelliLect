"""Programmatic fixtures for the extraction tests.

Everything here is built in-code from the same libraries the extractors use, so
the suite is fully self-contained — no binary fixtures checked into the repo and
nothing external to run.
"""

from __future__ import annotations

import io

import pymupdf
from docx import Document as open_docx
from PIL import Image
from pptx import Presentation
from pptx.util import Emu, Inches


def make_png(width: int = 24, height: int = 16, color: tuple[int, int, int] = (200, 30, 30)) -> bytes:
    """Return the bytes of a solid-color PNG of the given size."""
    buffer = io.BytesIO()
    Image.new("RGB", (width, height), color).save(buffer, format="PNG")
    return buffer.getvalue()


def make_docx() -> bytes:
    """A .docx with two heading levels, a paragraph, a table, and an embedded PNG.

    Structure (in body order):
        Heading 1: "Introduction"
        Heading 2: "Goals"
        Paragraph: "The goal is to extract text."
        Table: 2x2
        Picture: 24x16 PNG
    """
    document = open_docx()
    document.add_heading("Introduction", level=1)
    document.add_heading("Goals", level=2)
    document.add_paragraph("The goal is to extract text.")

    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Alpha"
    table.cell(1, 1).text = "1"

    document.add_picture(io.BytesIO(make_png()))

    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def make_pptx() -> bytes:
    """A .pptx with two slides: titles + body text, notes, and one image on slide 1."""
    presentation = Presentation()
    title_and_body = presentation.slide_layouts[1]

    slide1 = presentation.slides.add_slide(title_and_body)
    slide1.shapes.title.text = "First Slide"
    slide1.placeholders[1].text = "Body of the first slide."
    slide1.notes_slide.notes_text_frame.text = "Remember to explain the intro."
    slide1.shapes.add_picture(
        io.BytesIO(make_png(width=32, height=20, color=(20, 120, 200))),
        left=Inches(1),
        top=Inches(1),
        width=Emu(int(Inches(1))),
        height=Emu(int(Inches(0.625))),
    )

    slide2 = presentation.slides.add_slide(title_and_body)
    slide2.shapes.title.text = "Second Slide"
    slide2.placeholders[1].text = "Body of the second slide."

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def make_text_pdf() -> bytes:
    """A 2-page text PDF with a couple of positioned text blocks per page."""
    document = pymupdf.open()
    page1 = document.new_page()
    page1.insert_text((72, 100), "Chapter One")
    page1.insert_text((72, 160), "This is the first paragraph of the document.")

    page2 = document.new_page()
    page2.insert_text((72, 100), "Chapter Two")
    page2.insert_text((72, 160), "The second page continues the story.")

    data = document.tobytes()
    document.close()
    return data


def make_scanned_pdf(png_bytes: bytes | None = None) -> bytes:
    """A 1-page PDF that is a single full-page image with NO text layer."""
    if png_bytes is None:
        png_bytes = make_png(width=200, height=260, color=(90, 90, 90))
    document = pymupdf.open()
    page = document.new_page()
    page.insert_image(page.rect, stream=png_bytes)
    data = document.tobytes()
    document.close()
    return data


# MIME types the router recognizes.
PDF_CONTENT_TYPE = "application/pdf"
DOCX_CONTENT_TYPE = (
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
)
PPTX_CONTENT_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)
TXT_CONTENT_TYPE = "text/plain"
MARKDOWN_CONTENT_TYPE = "text/markdown"
